# Copyright (c) 2025 Kenneth Stott. MIT License.
# Canary: 97f8da42-9e90-48f0-b657-0e0641aab65f
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Unit tests for annotate() on XLSX, PDF, and PPTX extractors."""

from __future__ import annotations

import io

import pytest


def _make_xlsx_bytes(sheet_name: str = "Sheet1") -> bytes:
    pytest.importorskip("openpyxl")
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name
    ws.append(["Name", "Value"])
    for i in range(1, 6):
        ws.append([f"item{i}", i * 10])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pdf_bytes() -> bytes:
    pytest.importorskip("pypdf")
    try:
        from reportlab.pdfgen import canvas as rl_canvas
    except ImportError:
        pytest.skip("reportlab not installed; skipping PDF annotate test")

    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf)
    c.drawString(100, 750, "This is a test sentence on page one of the document.")
    c.showPage()
    c.drawString(100, 750, "This is another sentence on page two of the document.")
    c.showPage()
    c.save()
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    import pptx
    from pptx.util import Inches

    prs = pptx.Presentation()
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.name = "TextBox1"
    txBox.text_frame.text = "Hello from slide one of the presentation."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestXlsxAnnotate:
    def test_xlsx_annotate_stamps_sheet(self):
        from chonk.chunking import chunk_document
        from chonk.extractors._xlsx import XlsxExtractor

        ext = XlsxExtractor()
        data = _make_xlsx_bytes("MySheet")
        text = ext.extract(data)
        chunks = chunk_document("test", text, 10, 5000, 0.15, include_breadcrumb=False)
        annotated = ext.annotate(chunks, data)

        stamped = [c for c in annotated if c.source_detail is not None]
        assert stamped, "No chunks were stamped with source_detail"
        assert stamped[0].source_detail["sheet"] == "MySheet"


class TestPdfAnnotate:
    def test_pdf_annotate_stamps_page(self):
        from chonk.chunking import chunk_document
        from chonk.extractors._pdf import PdfExtractor

        ext = PdfExtractor()
        data = _make_pdf_bytes()
        text = ext.extract(data)
        chunks = chunk_document("test", text, 10, 5000, 0.15, include_breadcrumb=False)
        annotated = ext.annotate(chunks, data)

        stamped = [c for c in annotated if c.source_detail is not None]
        assert stamped, "No chunks were stamped with source_detail"
        detail = stamped[0].source_detail
        assert "page" in detail or "page_start" in detail


class TestPptxAnnotate:
    def test_pptx_annotate_stamps_slide(self):
        from chonk.chunking import chunk_document
        from chonk.extractors._pptx import PptxExtractor

        ext = PptxExtractor()
        data = _make_pptx_bytes()
        text = ext.extract(data)
        chunks = chunk_document("test", text, 10, 5000, 0.15, include_breadcrumb=False)
        annotated = ext.annotate(chunks, data)

        stamped = [c for c in annotated if c.source_detail is not None]
        assert stamped, "No chunks were stamped with source_detail"
        assert stamped[0].source_detail["slide"] == 1
