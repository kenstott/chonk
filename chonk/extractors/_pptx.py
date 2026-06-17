# Copyright (c) 2025 Kenneth Stott. MIT License.
# Canary: fb48752f-c49b-427f-a135-6c6a75f90ed4
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""PPTX text extractor using python-pptx."""

from __future__ import annotations

from io import BytesIO
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from chonk.models import DocumentChunk

try:
    import pptx as _pptx_module

    _PPTX_AVAILABLE = True
except ImportError:
    _pptx_module = None  # type: ignore[assignment]
    _PPTX_AVAILABLE = False


def _extract_pptx_content(prs: Any) -> str:  # noqa: ANN401
    """Extract text content from a python-pptx Presentation object."""
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []

        for shape in slide.shapes:
            shape_text_val: Any = getattr(shape, "text", None)
            if shape_text_val and str(shape_text_val).strip():
                slide_text.append(str(shape_text_val).strip())

            if shape.has_table:
                table_rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    slide_text.append("\n".join(table_rows))

        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))

    return "\n\n".join(slides)


class PptxExtractor:
    """Extract plain text from PPTX bytes."""

    def can_handle(self, doc_type: str) -> bool:
        return doc_type == "pptx"

    def extract(self, data: bytes, source_path: str | None = None) -> str:
        if not _PPTX_AVAILABLE:
            raise ImportError("pip install chonk[pptx]")
        assert _pptx_module is not None

        prs = _pptx_module.Presentation(BytesIO(data))
        return _extract_pptx_content(prs)

    def annotate(
        self,
        chunks: list[DocumentChunk],
        data: bytes,
        source_path: str | None = None,
    ) -> list[DocumentChunk]:
        if not _PPTX_AVAILABLE:
            raise ImportError("pip install chonk[pptx]")
        assert _pptx_module is not None

        prs = _pptx_module.Presentation(BytesIO(data))

        # Build per-slide shape index: slide_num → list of (shape_name_or_index, text)
        slide_shapes: list[tuple[int, str, str]] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                shape_label = shape.name if shape.name else str(shape_idx)
                shape_text_val: Any = getattr(shape, "text", None)
                if shape_text_val and str(shape_text_val).strip():
                    slide_shapes.append((slide_num, shape_label, str(shape_text_val).strip()))

        for chunk in chunks:
            content = chunk.content
            for slide_num, shape_label, shape_text in slide_shapes:
                if shape_text in content or any(
                    line.strip() in shape_text
                    for line in content.split("\n")
                    if len(line.strip()) > 10
                ):
                    chunk.source_detail = {"slide": slide_num, "shape": shape_label}
                    break

        return chunks
